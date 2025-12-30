from pptx import Presentation  # Import the Presentation class from the python-pptx library

"""
Extracts table data from a PowerPoint (.pptx) file.

This function reads a PowerPoint file and extracts data from tables present in the slides.
It processes each slide and each shape within the slides to find tables, then extracts the
headers and corresponding data from these tables. The data is returned as a list of dictionaries,
where each dictionary represents the data from one table.

Args:
    file_name (str): The path to the PowerPoint file from which to extract table data.

Returns:
    list: A list of dictionaries containing the extracted table data. Each dictionary represents
          the data from one table, with headers as keys and corresponding data as values.
"""


def extract_table_data_from_pptx(file_name):
    prs = Presentation(file_name)  # Open the PowerPoint file
    all_data = []  # Initialize a list to hold all extracted data

    for slide in prs.slides:  # Iterate through each slide in the presentation
        ctr=0
        all_data.append("-----------------------------")  # Add a separator for readability
        for shape in slide.shapes:  # Iterate through each shape in the slide
            # Append the slide title, append only the first text frame in the slide
            if shape.has_text_frame and ctr==0:
                ctr+=1
                all_data.append("Title: " + shape.text_frame.text)
            if shape.has_table:  # Check if the shape contains a table
                table = shape.table  # Get the table object
                slide_data = []  # Initialize a list to hold data for the current slide
                row_idx = 0  # Initialize row index
                while row_idx < len(table.rows):  # Iterate through the rows of the table
                    headers = [cell.text for cell in table.rows[row_idx].cells]  # Extract headers from the current row
                    if row_idx + 1 < len(table.rows):  # Check if there is a next row
                        data = [cell.text for cell in table.rows[row_idx + 1].cells]  # Extract data from the next row
                        slide_data.append(dict(zip(headers, data)))  # Combine headers and data into a dictionary and add to slide_data
                    row_idx += 2  # Move to the next pair of rows
                all_data.append(slide_data)  # Add the slide data to all_data

    return all_data  # Return the extracted data
def extract_table_data_from_pptx1(file_name):
    prs = Presentation(file_name)  # Open the PowerPoint file
    all_data = []  # Initialize a list to hold all extracted data

    for slide in prs.slides:  # Iterate through each slide in the presentation
        for shape in slide.shapes:  # Iterate through each shape in the slide
            if shape.has_table:  # Check if the shape contains a table
                all_data.append("######")  # Add a separator for readability
                table = shape.table  # Get the table object
                slide_data = []  # Initialize a list to hold data for the current slide
                row_idx = 0  # Initialize row index
                while row_idx < len(table.rows):  # Iterate through the rows of the table
                    headers = [cell.text for cell in table.rows[row_idx].cells]  # Extract headers from the current row
                    if row_idx + 1 < len(table.rows):  # Check if there is a next row
                        data = [cell.text for cell in table.rows[row_idx + 1].cells]  # Extract data from the next row
                        slide_data.append(dict(zip(headers, data)))  # Combine headers and data into a dictionary and add to slide_data
                    row_idx += 2  # Move to the next pair of rows
                all_data.append(slide_data)  # Add the slide data to all_data

    return all_data  # Return the extracted data

#extract_table_data_from_pptx(file_name)
file_name = '/Users/<PARAM_USER>/Downloads/a.pptx'
data = extract_table_data_from_pptx(file_name)

# format & print the data in a more readable way
for i in data:
    #if is string 
    if isinstance(i, str) and (i == "-----------------------------" or i.index("Title: ") >= 0):
        print(i)
    # if i is a non-empty list
    if isinstance(i, list) and len(i) > 0:
        for j in i:
            for k, v in j.items():
                if v:  # check if value is not empty
                    print(f"***{k}***\n{v}")
            print("\n")


